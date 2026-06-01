"""Data-driven rendering of the experience slides from the content store.

For each of the 15 categories the deck has a section divider (left untouched)
followed by one or more content slides holding a {client | description} table.
This module rebuilds those content slides from the store: it flows a category's
projects across as many slides as needed (paginating so nothing overflows),
cloning the first content slide as the template for any extra pages and deleting
surplus pages. Every generated slide keeps the original design.
"""

from __future__ import annotations

import re

from pptx.oxml.ns import qn

from era_agent.pptx_utils import (
    clone_slide, move_slide, delete_slide, slide_index, refill_table,
)
from era_agent.content import store

# General-description Section number -> category id (Sections 3–17).
_SEC2CAT = {
    3: "ma", 4: "competition", 5: "banking", 6: "treasury", 7: "grc",
    8: "sanctions", 9: "regulatory", 10: "public_policy", 11: "data_protection",
    12: "energy", 13: "distribution", 14: "contractual", 15: "hr_labour",
    16: "disputes", 17: "ip",
}
_SEC_RE = re.compile(r"Sec(?:t|ț|ţ)iunea?\s*(\d+)|Section\s*(\d+)")

# Pagination tuning. Calibrated to the template: the original deck's fullest
# experience slide held ~3,170 description chars over 3 rows and rendered with
# room to spare, so 3,200 is the proven-safe per-slide budget. _MAX_ROWS is a
# secondary guard against many-tiny-rows (the original packed up to 5 per slide).
_CHAR_BUDGET = 3200
_MAX_ROWS = 6


def _label_shape(slide):
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text
            if _SEC_RE.search(t) and len(t) < 90:
                return sh
    return None


def _table(slide):
    for sh in slide.shapes:
        if sh.has_table:
            return sh.table
    return None


def _mark_continuation(slide) -> None:
    sh = _label_shape(slide)
    if sh is None:
        return
    runs = sh.text_frame.paragraphs[0].runs
    if runs and "(cont.)" not in (runs[-1].text or ""):
        runs[-1].text = (runs[-1].text or "").rstrip() + " (cont.)"


def paginate(projects: list[tuple[str, str]],
             budget: int = _CHAR_BUDGET,
             max_rows: int = _MAX_ROWS) -> list[list[tuple[str, str]]]:
    """Split (client, desc) rows into pages by description-character budget."""
    pages: list[list[tuple[str, str]]] = []
    cur: list[tuple[str, str]] = []
    cur_len = 0
    for client, desc in projects:
        dlen = len(desc or "")
        if cur and (cur_len + dlen > budget or len(cur) >= max_rows):
            pages.append(cur)
            cur, cur_len = [], 0
        cur.append((client, desc))
        cur_len += dlen
    if cur:
        pages.append(cur)
    return pages


def build_projects_by_cat(lang: str) -> dict[str, list[tuple[str, str]]]:
    """Active projects grouped by category, as (client, long_desc) in `lang`."""
    out: dict[str, list[tuple[str, str]]] = {}
    for p in store.list_projects(active_only=True):
        client = (p["client"].get(lang) or p["client"].get("en") or "").strip()
        desc = (p["long"].get(lang) or p["long"].get("en") or "").strip()
        if not (client and desc):
            continue
        out.setdefault(p["category"], []).append((client, desc))
    return out


def _group_content_slides(prs) -> dict[str, list]:
    """Map category id -> ordered list of its content slides (table slides)."""
    groups: dict[str, list] = {}
    current = None
    for slide in prs.slides:
        table = None
        for sh in slide.shapes:
            if sh.has_text_frame:
                m = _SEC_RE.search(sh.text_frame.text)
                if m:
                    num = int(m.group(1) or m.group(2))
                    if num in _SEC2CAT:
                        current = _SEC2CAT[num]
            if sh.has_table:
                table = sh.table
        if table is not None and current is not None:
            groups.setdefault(current, []).append(slide)
    return groups


def _regen_category(prs, slides: list, pages: list[list[tuple[str, str]]]) -> None:
    proto = slides[0]
    page_slides = []
    for i in range(len(pages)):
        if i < len(slides):
            ps = slides[i]
        else:
            ps = clone_slide(prs, proto)
            move_slide(prs, ps, slide_index(prs, page_slides[-1]) + 1)
        page_slides.append(ps)

    # Drop surplus original pages (data shrank).
    for extra in slides[len(pages):]:
        delete_slide(prs, extra)

    for i, ps in enumerate(page_slides):
        tbl = _table(ps)
        if tbl is not None:
            refill_table(tbl, pages[i])
        if i > 0:
            _mark_continuation(ps)


def regenerate_experience(prs, lang: str,
                          budget: int = _CHAR_BUDGET,
                          max_rows: int = _MAX_ROWS) -> None:
    """Rebuild every experience section in `prs` from the content store."""
    projects_by_cat = build_projects_by_cat(lang)
    groups = _group_content_slides(prs)
    for cat, slides in groups.items():
        projects = projects_by_cat.get(cat, [])
        if not projects:
            continue  # no data: leave the template section as-is
        pages = paginate(projects, budget, max_rows)
        _regen_category(prs, slides, pages)
