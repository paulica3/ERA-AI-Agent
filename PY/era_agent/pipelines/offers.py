"""Custom Offer generation pipeline — fills ERA's custom_offer_{ro,en}.pptx template.

The template is a fully-formatted PowerPoint deck (18 slides). Generation only
*fills* the client-specific fields and reformats the fee text the user provides —
it never invents content, and crucially never invents pricing. Section 2 numbers
always come from the user; Claude only reformats that free text to match the
template's narrative style.
"""

import copy
import io
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

from era_agent.client import get_client
from era_agent.config import MODEL, MAX_TOKENS

TEMPLATES_DIR = Path(__file__).parent.parent.parent / "templates"

# Anchor strings used to locate the fields to fill in the template. These are the
# placeholder values baked into custom_offer_ro.pptx / custom_offer_en.pptx.
_CLIENT_NAME_PLACEHOLDER = "VITAFOR"
# The date string is baked in identically (Romanian) in both decks.
_DATE_PLACEHOLDER = "14 decembrie 2025"
_SALUTATION_ANCHOR = "Stimate domnule"  # RO salutation start; EN mirrors via index
_FEE_HEADING_RO = "Structura onorariilor"
_FEE_HEADING_EN = "Fee Structure"


# ── Low-level paragraph / run helpers ─────────────────────────────────────────

def _para_text(p) -> str:
    return "".join(r.text for r in p.runs)


def _set_para_text(p, text: str) -> None:
    """Collapse all runs into the first, preserving its formatting; drop <a:br>."""
    runs = p.runs
    # Remove explicit line breaks so we don't leave stray <a:br> behind.
    for br in p._p.findall(qn("a:br")):
        p._p.remove(br)
    if not runs:
        p.add_run().text = text
        return
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def _clone_para_with_text(template_p, text: str):
    """Deep-copy a paragraph element, then set it to a single run of `text`.

    Preserves the template paragraph's formatting (pPr + first run's rPr).
    """
    new_p = copy.deepcopy(template_p)
    # Strip explicit breaks.
    for br in new_p.findall(qn("a:br")):
        new_p.remove(br)
    runs = new_p.findall(qn("a:r"))
    if not runs:
        return new_p
    # Set first run's text, remove the rest.
    t = runs[0].find(qn("a:t"))
    if t is None:
        t = runs[0].makeelement(qn("a:t"), {})
        runs[0].append(t)
    t.text = text
    for r in runs[1:]:
        new_p.remove(r)
    return new_p


def _find_shape(slide, anchor: str):
    """Return (shape, paragraph_index) of the first text shape containing anchor."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for i, p in enumerate(sh.text_frame.paragraphs):
            if anchor in _para_text(p):
                return sh, i
    return None, None


# ── Fee section (Section 2) ───────────────────────────────────────────────────

def reformat_fee_structure(fee_text: str, lang: str = "ro") -> list[str]:
    """Reformat the user's free-text fee description into polished paragraphs
    matching the template's Section 2 style.

    HARD CONSTRAINT: never invent, alter, round or add any numbers, amounts,
    rates, hours, currencies or percentages. Every figure in the output must
    appear verbatim in the input. Claude only rephrases / structures the prose.
    """
    fee_text = (fee_text or "").strip()
    if not fee_text:
        return []

    lang_name = "română" if lang == "ro" else "engleză"
    system = (
        "Ești asistentul juridic al firmei Efrim, Roșca & Asociații din Republica Moldova. "
        "Reformatezi textul despre onorarii furnizat de avocat într-un stil profesionist, "
        "potrivit pentru secțiunea financiară a unei oferte de servicii juridice."
    )
    prompt = (
        f"Reformatează textul de mai jos despre structura onorariilor într-un set de "
        f"paragrafe profesioniste, în limba {lang_name}, în stilul unei oferte juridice.\n\n"
        "REGULI ABSOLUTE:\n"
        "- NU inventa, NU modifica, NU rotunji și NU adăuga NICIUN număr, sumă, tarif, "
        "număr de ore, monedă sau procent. Fiecare cifră din răspuns TREBUIE să apară "
        "exact în textul de intrare.\n"
        "- Nu adăuga informații noi care nu sunt în text. Doar reorganizează și "
        "rafinează formularea.\n"
        "- Păstrează toate sumele și intervalele exact cum sunt scrise.\n\n"
        "Returnează STRICT un array JSON de string-uri (un string per paragraf), "
        "fără text înainte sau după.\n\n"
        f"Text:\n{fee_text}"
    )

    client = get_client()
    resp = client.messages.create(
        model=MODEL,
        max_tokens=MAX_TOKENS,
        system=system,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = "\n".join(b.text for b in resp.content if hasattr(b, "text")).strip()
    clean = re.sub(r"^```(?:json)?\s*|\s*```$", "", raw, flags=re.MULTILINE).strip()
    try:
        data = json.loads(clean)
    except json.JSONDecodeError:
        # Fall back to splitting on blank lines if Claude didn't return JSON.
        return [blk.strip() for blk in re.split(r"\n\s*\n", fee_text) if blk.strip()]
    if isinstance(data, list):
        return [str(x).strip() for x in data if str(x).strip()]
    return [str(data).strip()]


def compose_cover_letter(context: str, client_name: str, lang: str = "ro") -> list[str]:
    """Compose the slide-2 cover-letter body from the user's free context.

    The user describes the matter (who the client is, what the transaction /
    mandate is, what assistance is needed). Claude turns that into the formal
    body of an ERA offer letter — exactly 3 paragraphs matching the house style:
      1. ERA is pleased to submit this offer in response to the client's request
         regarding <the matter>.
      2. The scope of assistance, from the perspective of Moldovan law.
      3. The standard closing (attached: scope, team, assumptions, fee estimate;
         remaining available for clarifications).

    The salutation and the "Cu respect, / Oleg EFRIM / Managing Partner" block are
    handled separately, so this returns ONLY the body paragraphs. Claude must not
    invent specific facts (figures, names, dates) that aren't in the context.
    """
    context = (context or "").strip()
    if not context:
        return []

    if lang == "ro":
        system = (
            "Ești asistentul juridic al firmei de avocatură „Efrim, Roșca & Asociații” "
            "din Republica Moldova. Redactezi corpul scrisorii de însoțire a unei oferte "
            "de servicii juridice, într-un registru formal, sobru și profesionist."
        )
        prompt = (
            "Pe baza contextului de mai jos (descris liber de avocat), redactează corpul "
            "scrisorii de însoțire a ofertei de servicii juridice.\n\n"
            "Structură — EXACT 3 paragrafe, în această ordine:\n"
            "1. „Efrim, Roșca & Asociații are plăcerea de a vă transmite această ofertă de "
            "servicii juridice, ca răspuns la solicitarea privind …” — descrie pe scurt "
            "obiectul mandatului/tranzacției din context.\n"
            "2. Domeniul asistenței, din perspectiva dreptului Republicii Moldova "
            "(etapele esențiale relevante pentru acest tip de mandat).\n"
            "3. Paragraf de încheiere standard: „Anexat veți regăsi domeniul propus al "
            "serviciilor, componența echipei, ipotezele de lucru și estimarea onorariilor. "
            "Rămânem la dispoziția dumneavoastră pentru orice clarificări …”.\n\n"
            "REGULI:\n"
            "- NU inventa nume, cifre, sume, date sau fapte specifice care nu apar în context.\n"
            "- Folosește un ton formal, diacritice corecte, fără titluri sau bullet-uri.\n"
            "- Returnează STRICT un array JSON cu exact 3 string-uri (un paragraf fiecare), "
            "fără text înainte sau după.\n\n"
            f"Numele clientului: {client_name}\n\n"
            f"Context:\n{context}"
        )
    else:
        system = (
            "You are the legal assistant of the law firm “Efrim, Roșca & Asociații” in the "
            "Republic of Moldova. You draft the body of the cover letter accompanying an "
            "offer of legal services, in a formal, restrained and professional register."
        )
        prompt = (
            "Based on the context below (described freely by the lawyer), draft the body of "
            "the cover letter accompanying an offer of legal services.\n\n"
            "Structure — EXACTLY 3 paragraphs, in this order:\n"
            "1. “Efrim, Roșca & Asociații is pleased to submit this offer of legal services, "
            "in response to the request regarding …” — briefly describe the subject of the "
            "mandate/transaction from the context.\n"
            "2. The scope of assistance, from the perspective of the law of the Republic of "
            "Moldova (the essential stages relevant to this type of mandate).\n"
            "3. A standard closing paragraph: “Attached you will find the proposed scope of "
            "services, the team composition, the working assumptions and the fee estimate. "
            "We remain at your disposal for any clarifications …”.\n\n"
            "RULES:\n"
            "- Do NOT invent names, figures, amounts, dates or specific facts not present "
            "in the context.\n"
            "- Use a formal tone, no headings or bullet points.\n"
            "- Return STRICTLY a JSON array of exactly 3 strings (one paragraph each), "
            "with no text before or after.\n\n"
            f"Client name: {client_name}\n\n"
            f"Context:\n{context}"
        )

    client = get_client()
    resp = client.messages.create(
        model=MODEL,
        max_tokens=MAX_TOKENS,
        system=system,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = "\n".join(b.text for b in resp.content if hasattr(b, "text")).strip()
    clean = re.sub(r"^```(?:json)?\s*|\s*```$", "", raw, flags=re.MULTILINE).strip()
    try:
        data = json.loads(clean)
    except json.JSONDecodeError:
        # Fall back to splitting the raw context on blank lines.
        return [blk.strip() for blk in re.split(r"\n\s*\n", context) if blk.strip()][:3]
    if isinstance(data, list):
        return [str(x).strip() for x in data if str(x).strip()][:3]
    return [str(data).strip()]


def _fill_fee_section(shape, paragraphs: list[str], lang: str) -> None:
    """Replace the fee-narrative body (everything after the heading) with new paras."""
    if not paragraphs:
        return
    heading = _FEE_HEADING_RO if lang == "ro" else _FEE_HEADING_EN
    tx = shape.text_frame._txBody
    p_els = tx.findall(qn("a:p"))

    # Locate the heading paragraph; body = everything after it.
    hi = None
    for i, pe in enumerate(p_els):
        txt = "".join(t.text or "" for t in pe.findall(f".//{qn('a:t')}"))
        if heading in txt:
            hi = i
            break
    if hi is None or hi + 1 >= len(p_els):
        return

    template_p = p_els[hi + 1]
    # Remove existing body paragraphs.
    for pe in p_els[hi + 1:]:
        tx.remove(pe)
    # Append rebuilt body paragraphs cloned from the template.
    for text in paragraphs:
        tx.append(_clone_para_with_text(template_p, text))


# ── Main entry point ──────────────────────────────────────────────────────────

def generate_custom_offer(
    client_name: str,
    date: str,
    addressee_salutation: str,
    addressee_block: str,
    intro_text: str = "",
    compose_intro: bool = True,
    fee_text: str = "",
    signatory_name: str = "Oleg EFRIM",
    signatory_title: str = "Managing Partner",
    lang: str = "ro",
    reformat_fees: bool = True,
) -> bytes:
    """Fill the Custom Offer template and return the .pptx bytes.

    Parameters
    ----------
    client_name : client/company name (slide 1 hero + addressee block reference).
    date : offer date string already formatted in the target language.
    addressee_salutation : e.g. "Stimate domnule BÎRCĂ" / "Dear Mr BÎRCĂ".
    addressee_block : multi-line addressee text for the slide-2 table (name,
        title, company, address — one per line).
    intro_text : the cover-letter content. If `compose_intro` is True this is the
        free-form *context* of the matter, which Claude turns into the formal
        letter body; otherwise it is used as literal paragraphs (split on blank
        lines). Unused body slots in the template are blanked.
    compose_intro : if True, Claude composes the letter body from `intro_text`.
    fee_text : user-provided fee description (the numbers always come from here).
    reformat_fees : if True, Claude reformats `fee_text`; if False it is used
        verbatim (split on blank lines).
    lang : "ro" or "en" — picks the matching template.
    """
    # Resolve the cover-letter body paragraphs.
    if intro_text.strip():
        if compose_intro:
            intro_paragraphs = compose_cover_letter(intro_text, client_name, lang=lang)
        else:
            intro_paragraphs = [
                b.strip() for b in re.split(r"\n\s*\n", intro_text) if b.strip()
            ]
    else:
        intro_paragraphs = []

    tmpl = "custom_offer_ro.pptx" if lang == "ro" else "custom_offer_en.pptx"
    prs = Presentation(TEMPLATES_DIR / tmpl)

    # ── Slide 1: client name + date hero ──────────────────────────────
    s1 = prs.slides[0]
    for sh in s1.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            t = _para_text(p)
            if t.strip() == _CLIENT_NAME_PLACEHOLDER:
                _set_para_text(p, client_name)
            elif _DATE_PLACEHOLDER in t:
                _set_para_text(p, date)

    # ── Slide 2: addressee table + salutation + cover-letter body ─────
    s2 = prs.slides[1]

    # Addressee table cell (shape with a table; r0c1 holds the addressee block).
    for sh in s2.shapes:
        if sh.has_table:
            cell = sh.table.rows[0].cells[1]
            # Indent style in template uses leading spaces; mirror by setting the
            # cell's first paragraph and clearing the rest.
            lines = [ln for ln in addressee_block.splitlines()]
            cps = cell.text_frame.paragraphs
            for idx, p in enumerate(cps):
                if idx < len(lines):
                    _set_para_text(p, lines[idx])
                else:
                    _set_para_text(p, "")
            # If more lines than existing paragraphs, append clones of the first.
            if len(lines) > len(cps):
                tx = cell.text_frame._txBody
                template_p = tx.findall(qn("a:p"))[0]
                for extra in lines[len(cps):]:
                    tx.append(_clone_para_with_text(template_p, extra))
            break

    # Salutation + body + signatory live in the letter text shape.
    sal_anchor = _SALUTATION_ANCHOR if lang == "ro" else "Dear"
    letter_sh, _ = _find_shape(s2, sal_anchor)
    if letter_sh is None:
        # EN salutation may differ; fall back to the shape holding the signatory.
        letter_sh, _ = _find_shape(s2, signatory_name)
    if letter_sh is not None:
        # Non-empty paragraphs appear in this fixed order:
        #   [salutation, body1, body2, body3, closing, name, title]
        _CLOSINGS = {
            "Cu respect,", "Yours sincerely,", "Sincerely,",
            "Kind regards,", "Respectfully,",
        }
        non_empty = [p for p in letter_sh.text_frame.paragraphs if _para_text(p).strip()]
        filled = (intro_paragraphs or [])[:3]

        # Find the closing paragraph; it separates body from signatory block.
        closing_idx = next(
            (i for i, p in enumerate(non_empty) if _para_text(p).strip() in _CLOSINGS),
            None,
        )

        if non_empty:
            _set_para_text(non_empty[0], addressee_salutation)  # salutation
        # Body paragraphs sit between salutation and closing.
        body_end = closing_idx if closing_idx is not None else len(non_empty)
        body_paras = non_empty[1:body_end]
        for slot, p in enumerate(body_paras):
            # Fill from composed paragraphs; blank any leftover template slots so
            # no placeholder prose (e.g. the VITAFOR text) ever leaks through.
            _set_para_text(p, filled[slot] if slot < len(filled) else "")
        # Signatory block: the two paragraphs after the closing.
        if closing_idx is not None:
            sig = non_empty[closing_idx + 1:]
            if len(sig) >= 1:
                _set_para_text(sig[0], signatory_name)
            if len(sig) >= 2:
                _set_para_text(sig[1], signatory_title)

    # ── Slide 9: Section 2 fee structure ──────────────────────────────
    if fee_text.strip():
        if reformat_fees:
            fee_paras = reformat_fee_structure(fee_text, lang=lang)
        else:
            fee_paras = [b.strip() for b in re.split(r"\n\s*\n", fee_text) if b.strip()]
        heading = _FEE_HEADING_RO if lang == "ro" else _FEE_HEADING_EN
        s9 = prs.slides[8]
        fee_sh, _ = _find_shape(s9, heading)
        if fee_sh is not None:
            _fill_fee_section(fee_sh, fee_paras, lang)

    # ── Serialise ─────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
