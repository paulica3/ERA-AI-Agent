"""General Description generation pipeline — fills the firm's standing
credentials deck (general_description_{ro,en}.pptx).

Unlike the Custom Offer, this deck is a fixed firm-overview brochure. Generation
only fills a handful of per-client fields and never touches the 15 "relevant
experience" sections:

  • Slide 1  — offer date.
  • Slide 2  — addressee block (To: / Către:), salutation, an OPTIONAL
               AI-composed opening paragraph tailored from free-text context,
               and the signatory name/title.
  • Slide 12 — the standard hourly rate (defaults to EUR 250 when left blank).

The AI step (compose_intro_paragraph) only rephrases the lawyer's free-text
context into one formal opening paragraph; it never invents names, numbers,
dates or facts that are not in the context.
"""

import copy
import datetime
import io
import re
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

from era_agent.client import get_client
from era_agent.config import MODEL, MAX_TOKENS
from era_agent.pipelines.experience import regenerate_experience

TEMPLATES_DIR = Path(__file__).parent.parent.parent / "templates"

_DATE_PLACEHOLDER = "27 May 2026"
_RATE_ANCHOR = "EUR 250"

_RO_MONTHS = {
    1: "ianuarie", 2: "februarie", 3: "martie", 4: "aprilie", 5: "mai",
    6: "iunie", 7: "iulie", 8: "august", 9: "septembrie", 10: "octombrie",
    11: "noiembrie", 12: "decembrie",
}


# ── Low-level paragraph/run helpers ───────────────────────────────────────────

def _para_text(p) -> str:
    return "".join(r.text for r in p.runs)


def _set_para_text(p, text: str) -> None:
    """Collapse all runs into the first, preserving its formatting; drop <a:br>."""
    for br in p._p.findall(qn("a:br")):
        p._p.remove(br)
    runs = p.runs
    if not runs:
        p.add_run().text = text
        return
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def _find_shape_containing(slide, substr: str):
    """Return the first text-frame shape whose text contains substr."""
    for sh in slide.shapes:
        if sh.has_text_frame and substr in sh.text_frame.text:
            return sh
    return None


def _set_first_para_matching(shape, predicate, text: str) -> bool:
    """Set the text of the first paragraph in shape matching predicate(text)."""
    if shape is None:
        return False
    for p in shape.text_frame.paragraphs:
        if predicate(_para_text(p).strip()):
            _set_para_text(p, text)
            return True
    return False


# ── Slide 2: addressee table cell ─────────────────────────────────────────────

def _set_addressee_cell(table, lines: list[str], lang: str) -> None:
    """Rebuild the addressee cell (row 0, col 1) as a single paragraph:
    'To:  <line0>' (bold) followed by line breaks + the remaining lines.
    """
    lines = [ln for ln in (lines or []) if ln is not None]
    if not lines:
        return
    label = "To:  " if lang == "en" else "Către:  "
    cell = table.rows[0].cells[1]
    txBody = cell.text_frame._txBody
    ps = txBody.findall(qn("a:p"))
    if not ps:
        return
    first_p = ps[0]

    # Capture bold and normal run-property templates from the existing runs.
    bold_rpr = norm_rpr = None
    for r in first_p.findall(qn("a:r")):
        rpr = r.find(qn("a:rPr"))
        if rpr is None:
            continue
        if rpr.get("b") == "1" and bold_rpr is None:
            bold_rpr = copy.deepcopy(rpr)
        elif rpr.get("b") in (None, "0") and norm_rpr is None:
            norm_rpr = copy.deepcopy(rpr)
    if bold_rpr is None and first_p.findall(qn("a:r")):
        bold_rpr = copy.deepcopy(first_p.findall(qn("a:r"))[0].find(qn("a:rPr")))
    if norm_rpr is None:
        norm_rpr = copy.deepcopy(bold_rpr) if bold_rpr is not None else None

    # Clear existing runs / breaks / endParaRPr (keep pPr).
    for tag in ("a:r", "a:br", "a:endParaRPr"):
        for el in first_p.findall(qn(tag)):
            first_p.remove(el)

    def mk_run(text: str, rpr):
        r = first_p.makeelement(qn("a:r"), {})
        if rpr is not None:
            r.append(copy.deepcopy(rpr))
        t = first_p.makeelement(qn("a:t"), {})
        t.text = text
        r.append(t)
        return r

    first_p.append(mk_run(label + lines[0], bold_rpr))
    for ln in lines[1:]:
        first_p.append(first_p.makeelement(qn("a:br"), {}))
        first_p.append(mk_run(ln, norm_rpr))

    # Blank any extra paragraphs in the cell so no stale lines remain.
    for extra_p in ps[1:]:
        for tag in ("a:r", "a:br"):
            for el in extra_p.findall(qn(tag)):
                extra_p.remove(el)


# ── Slide 2: optional AI-composed opening paragraph ───────────────────────────

def compose_intro_paragraph(context: str, lang: str = "ro") -> str:
    """Turn the lawyer's free-text context into ONE formal opening paragraph
    that leads into the standard firm-overview letter. Never invents facts."""
    context = (context or "").strip()
    if not context:
        return ""

    if lang == "ro":
        system = (
            "Ești asistentul juridic al firmei „Efrim, Roșca & Asociații” din "
            "Republica Moldova. Redactezi un singur paragraf introductiv, formal "
            "și sobru, pentru scrisoarea de însoțire a unei prezentări generale a "
            "serviciilor firmei."
        )
        prompt = (
            "Pe baza contextului de mai jos, redactează UN SINGUR paragraf "
            "introductiv (2–3 propoziții) care personalizează scrisoarea pentru "
            "client și conduce spre prezentarea generală a serviciilor firmei.\n\n"
            "REGULI:\n"
            "- NU inventa nume, cifre, sume, date sau fapte care nu apar în context.\n"
            "- Ton formal, diacritice corecte, fără titluri sau bullet-uri.\n"
            "- Returnează DOAR textul paragrafului, fără ghilimele sau alt text.\n\n"
            f"Context:\n{context}"
        )
    else:
        system = (
            "You are the legal assistant of the law firm “Efrim, Roșca & "
            "Asociații” in the Republic of Moldova. You draft a single formal, "
            "restrained opening paragraph for the cover letter of a general "
            "description of the firm's services."
        )
        prompt = (
            "Based on the context below, draft ONE opening paragraph (2–3 "
            "sentences) that personalises the letter for the client and leads "
            "into the firm's general service overview.\n\n"
            "RULES:\n"
            "- Do NOT invent names, figures, amounts, dates or facts not in the context.\n"
            "- Formal tone, no headings or bullet points.\n"
            "- Return ONLY the paragraph text, with no quotes or extra text.\n\n"
            f"Context:\n{context}"
        )

    client = get_client()
    resp = client.messages.create(
        model=MODEL, max_tokens=MAX_TOKENS, system=system,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = "\n".join(b.text for b in resp.content if hasattr(b, "text")).strip()
    raw = re.sub(r"^```.*?\n|```$", "", raw, flags=re.DOTALL).strip()
    return raw.strip('"').strip()


def _insert_intro_paragraph(shape, text: str, lang: str) -> None:
    """Insert a new paragraph (cloned from the first standard body paragraph)
    holding the composed intro, right before that standard paragraph."""
    if shape is None or not text:
        return
    lead_anchor = "is pleased to provide" if lang == "en" else "are plăcerea"
    target_p = None
    for p in shape.text_frame.paragraphs:
        if lead_anchor in _para_text(p):
            target_p = p
            break
    if target_p is None:
        return
    new_p = copy.deepcopy(target_p._p)
    # Set the clone to a single run of `text`.
    for br in new_p.findall(qn("a:br")):
        new_p.remove(br)
    runs = new_p.findall(qn("a:r"))
    if runs:
        t = runs[0].find(qn("a:t"))
        if t is None:
            t = runs[0].makeelement(qn("a:t"), {})
            runs[0].append(t)
        t.text = text
        for r in runs[1:]:
            new_p.remove(r)
    target_p._p.addprevious(new_p)


# ── Slide 12: hourly rate ─────────────────────────────────────────────────────

def _set_hourly_rate(slide, rate: str) -> None:
    rate = (rate or "").strip()
    if not rate:
        return  # leave the template's default EUR 250
    replacement = rate if rate.upper().startswith("EUR") else f"EUR {rate}"
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if _RATE_ANCHOR in (r.text or ""):
                    r.text = r.text.replace(_RATE_ANCHOR, replacement)
                    return


# ── Main entry point ──────────────────────────────────────────────────────────

def generate_general_description(
    addressee_block: str,
    addressee_salutation: str,
    date: str = "",
    intro_context: str = "",
    compose_intro: bool = True,
    signatory_name: str = "Oleg EFRIM",
    signatory_title: str = "Managing Partner",
    hourly_rate: str = "",
    lang: str = "ro",
) -> bytes:
    """Fill the General Description template and return the .pptx bytes.

    Parameters
    ----------
    addressee_block : multi-line addressee text (name, title, company, …); the
        first line is rendered in bold after the 'To:' / 'Către:' label.
    addressee_salutation : e.g. "Dear Mr DUMBRAVĂ," / "Stimate Domnule DUMBRAVĂ,".
    date : offer date already formatted; if empty, today's date in `lang`.
    intro_context : optional free-text matter context. If `compose_intro` is
        True, Claude turns it into one formal opening paragraph inserted at the
        top of the letter body.
    hourly_rate : optional rate; if empty the template keeps EUR 250.
    lang : "ro" or "en" — picks the matching template.
    """
    lang = "en" if lang == "en" else "ro"
    tmpl = f"general_description_{lang}.pptx"
    prs = Presentation(TEMPLATES_DIR / tmpl)

    # Resolve date default.
    if not date.strip():
        today = datetime.date.today()
        if lang == "ro":
            date = f"{today.day} {_RO_MONTHS[today.month]} {today.year}"
        else:
            date = f"{today.day} {today.strftime('%B')} {today.year}"

    # ── Slide 1: date ─────────────────────────────────────────────────
    s1 = prs.slides[0]
    _set_first_para_matching(
        _find_shape_containing(s1, _DATE_PLACEHOLDER),
        lambda t: t == _DATE_PLACEHOLDER,
        date,
    )

    # ── Slide 2: addressee, salutation, intro, signatory ──────────────
    s2 = prs.slides[1]

    for sh in s2.shapes:
        if sh.has_table:
            lines = [ln for ln in addressee_block.splitlines() if ln.strip()]
            _set_addressee_cell(sh.table, lines, lang)
            break

    letter = _find_shape_containing(s2, "Yours sincerely") \
        or _find_shape_containing(s2, "Al dumneavoastră") \
        or _find_shape_containing(s2, signatory_name)

    if letter is not None:
        sal_prefix = "Dear" if lang == "en" else "Stimate"
        _set_first_para_matching(
            letter, lambda t: t.startswith(sal_prefix), addressee_salutation
        )
        # Signatory name (keep trailing comma to match template style).
        name_txt = signatory_name if signatory_name.rstrip().endswith(",") \
            else f"{signatory_name},"
        _set_first_para_matching(
            letter, lambda t: "EFRIM" in t.upper() and "," in t, name_txt
        )
        _set_first_para_matching(
            letter,
            lambda t: t in ("Managing Partner", "Partener Managing",
                            "Partener Coordonator", "Managing  Associate"),
            signatory_title,
        )
        # Optional AI-composed opening paragraph.
        if intro_context.strip() and compose_intro:
            intro = compose_intro_paragraph(intro_context, lang=lang)
            _insert_intro_paragraph(letter, intro, lang)
        elif intro_context.strip():
            _insert_intro_paragraph(letter, intro_context.strip(), lang)

    # ── Slide 12: hourly rate ─────────────────────────────────────────
    _set_hourly_rate(prs.slides[11], hourly_rate)

    # ── Experience sections: rebuild from the content store ───────────
    # Done last so the fixed-index fills above (slides 1/2/12) are unaffected
    # by the slides this adds/removes further down the deck.
    regenerate_experience(prs, lang)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
