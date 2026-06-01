"""One-time migration: extract the firm's existing projects from the two decks
into the content store (projects.json).

  • Long descriptions + client names come from general_description_{en,ro}.pptx
    (all 15 categories).
  • Short descriptions come from custom_offer_{en,ro}.pptx (M&A slide only,
    matched to the same clients) — that is the only place short text exists today.

The EN and RO templates share an identical structure, so we extract each in the
same order and zip them to pair the two languages.

Run:  python -m era_agent.migration.seed_projects [--write]
Without --write it prints a summary (dry run).
"""

from __future__ import annotations

import re
import sys
import unicodedata
from pathlib import Path

from pptx import Presentation

from era_agent.content.schema import Project
from era_agent.content import store

TEMPLATES = Path(__file__).resolve().parent.parent.parent / "templates"

# General-description Section number -> category id.
SEC2CAT = {
    3: "ma", 4: "competition", 5: "banking", 6: "treasury", 7: "grc",
    8: "sanctions", 9: "regulatory", 10: "public_policy", 11: "data_protection",
    12: "energy", 13: "distribution", 14: "contractual", 15: "hr_labour",
    16: "disputes", 17: "ip",
}
_SEC_RE = re.compile(r"Sec(?:t|ț|ţ)iunea?\s*(\d+)|Section\s*(\d+)")


def _clean(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "").replace("\x0b", " ")).strip()


def _norm(s: str) -> str:
    """Normalise a client name for cross-deck matching."""
    s = unicodedata.normalize("NFKD", s or "")
    s = "".join(c for c in s if not unicodedata.combining(c))
    return re.sub(r"[^a-z0-9]+", " ", s.lower()).strip()


def _extract_general(lang: str) -> list[tuple[str, str, str]]:
    """Return ordered (category_id, client, long_desc) from the GD deck."""
    prs = Presentation(TEMPLATES / f"general_description_{lang}.pptx")
    out: list[tuple[str, str, str]] = []
    current = None
    for slide in prs.slides:
        table = None
        for sh in slide.shapes:
            if sh.has_text_frame:
                m = _SEC_RE.search(sh.text_frame.text)
                if m:
                    num = int(m.group(1) or m.group(2))
                    if num in SEC2CAT:
                        current = SEC2CAT[num]
            if sh.has_table:
                table = sh.table
        if table is None or current is None:
            continue
        for row in table.rows:
            client = _clean(row.cells[0].text)
            desc = _clean(row.cells[1].text) if len(row.cells) > 1 else ""
            if client and desc:
                out.append((current, client, desc))
    return out


def _extract_offer_ma(lang: str) -> dict[str, str]:
    """Return {normalised client -> short_desc} from the Custom Offer M&A slide."""
    prs = Presentation(TEMPLATES / f"custom_offer_{lang}.pptx")
    result: dict[str, str] = {}
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_table:
                continue
            for row in sh.table.rows:
                client = _clean(row.cells[0].text)
                desc = _clean(row.cells[1].text) if len(row.cells) > 1 else ""
                if client and desc:
                    result[_norm(client)] = desc
    return result


def build_projects() -> list[dict]:
    gen_en = _extract_general("en")
    gen_ro = _extract_general("ro")
    if len(gen_en) != len(gen_ro):
        print(f"  ! EN/RO row count differs ({len(gen_en)} vs {len(gen_ro)}); "
              f"pairing the first {min(len(gen_en), len(gen_ro))}.", file=sys.stderr)

    short_en = _extract_offer_ma("en")
    short_ro = _extract_offer_ma("ro")

    projects: list[dict] = []
    order_by_cat: dict[str, int] = {}
    for (cat, client_en, long_en), (_, client_ro, long_ro) in zip(gen_en, gen_ro):
        order_by_cat[cat] = order_by_cat.get(cat, 0) + 1
        key = _norm(client_en)
        p = Project(
            category=cat,
            client={"en": client_en, "ro": client_ro},
            long={"en": long_en, "ro": long_ro},
            short={"en": short_en.get(key, ""), "ro": short_ro.get(key, "")},
            order=order_by_cat[cat],
        )
        projects.append(p.to_dict())
    return projects


def main() -> None:
    write = "--write" in sys.argv
    projects = build_projects()

    by_cat: dict[str, int] = {}
    with_short = 0
    for p in projects:
        by_cat[p["category"]] = by_cat.get(p["category"], 0) + 1
        if p["short"]["en"] or p["short"]["ro"]:
            with_short += 1

    print(f"Extracted {len(projects)} projects ({with_short} with a short version):")
    for cat, n in by_cat.items():
        print(f"  {cat:16s} {n}")

    if write:
        import json
        from era_agent.content.schema import CATEGORIES
        payload = {"version": 1, "categories": CATEGORIES, "projects": projects}
        store.SEED_PATH.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        print(f"\nWrote {len(projects)} projects to {store.SEED_PATH}")
    else:
        print("\n(dry run — pass --write to persist)")


if __name__ == "__main__":
    main()
