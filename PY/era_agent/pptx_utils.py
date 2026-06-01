"""Low-level PPTX helpers for data-driven slide generation.

The decks are built by *cloning* an existing template slide (so every generated
slide keeps the exact design — fonts, footer, logo, section label) and then
refilling its table. python-pptx has no public slide-clone, so we copy the
source slide's shape tree into a fresh slide bound to the same layout, and
re-create the image relationships with remapped rIds.
"""

from __future__ import annotations

import copy

from pptx.oxml.ns import qn


def clone_slide(prs, src_slide):
    """Append a deep copy of src_slide to prs and return the new slide.

    Copies all shapes and re-creates non-layout relationships (e.g. images),
    remapping their rIds so blip references stay valid in the new slide part.
    """
    new_slide = prs.slides.add_slide(src_slide.slide_layout)

    # Remove placeholders that add_slide injected from the layout.
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)

    # Deep-copy every shape element from the source.
    spTree = new_slide.shapes._spTree
    for shp in src_slide.shapes:
        spTree.append(copy.deepcopy(shp._element))

    # Re-create relationships (images, hyperlinks, …) and remap rIds.
    rid_map: dict[str, str] = {}
    for rid, rel in src_slide.part.rels.items():
        reltype = rel.reltype
        if reltype.endswith("/slideLayout") or reltype.endswith("/notesSlide"):
            continue
        if rel.is_external:
            new_rid = new_slide.part.relate_to(rel.target_ref, reltype, is_external=True)
        else:
            new_rid = new_slide.part.relate_to(rel.target_part, reltype)
        rid_map[rid] = new_rid

    if rid_map:
        _remap_rids(spTree, rid_map)

    return new_slide


def _remap_rids(element, rid_map: dict[str, str]) -> None:
    """Rewrite r:embed / r:id / r:link attributes in a copied subtree."""
    attrs = (qn("r:embed"), qn("r:id"), qn("r:link"))
    for el in element.iter():
        for attr in attrs:
            val = el.get(attr)
            if val in rid_map:
                el.set(attr, rid_map[val])


def move_slide(prs, slide, new_index: int) -> None:
    """Move an already-added slide to position new_index in the deck order."""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    # Find the sldId whose r:id resolves to this slide's part.
    target_rid = None
    for rid, part in prs.part.rels.items():
        if getattr(part, "target_part", None) is slide.part:
            target_rid = rid
            break
    if target_rid is None:
        # rels iterate gives (rid, rel); fall back via rel objects
        for rid, rel in prs.part.rels.items():
            if not rel.is_external and rel.target_part is slide.part:
                target_rid = rid
                break
    for sld in ids:
        if sld.get(qn("r:id")) == target_rid:
            sldIdLst.remove(sld)
            sldIdLst.insert(new_index, sld)
            return


def delete_slide(prs, slide) -> None:
    """Remove a slide from the deck (drops it from sldIdLst)."""
    sldIdLst = prs.slides._sldIdLst
    for sld in list(sldIdLst):
        rid = sld.get(qn("r:id"))
        rel = prs.part.rels.get(rid)
        if rel is not None and not rel.is_external and rel.target_part is slide.part:
            sldIdLst.remove(sld)
            return


def slide_index(prs, slide) -> int:
    """Return the current position of slide in the deck order."""
    for i, s in enumerate(prs.slides):
        if s is slide:
            return i
    return -1


# ── Table helpers ─────────────────────────────────────────────────────────────

def set_cell_text(cell, text: str) -> None:
    """Set a table cell's text, preserving its first run's formatting (rPr) and
    first paragraph's properties (pPr). Collapses to a single paragraph/run."""
    txBody = cell._tc.find(qn("a:txBody"))
    ps = txBody.findall(qn("a:p"))
    first = ps[0]
    for p in ps[1:]:
        txBody.remove(p)
    for br in first.findall(qn("a:br")):
        first.remove(br)
    runs = first.findall(qn("a:r"))
    if runs:
        t = runs[0].find(qn("a:t"))
        if t is None:
            t = runs[0].makeelement(qn("a:t"), {})
            runs[0].append(t)
        t.text = text
        for r in runs[1:]:
            first.remove(r)
    else:
        r = first.makeelement(qn("a:r"), {})
        t = first.makeelement(qn("a:t"), {})
        t.text = text
        r.append(t)
        first.append(r)


def refill_table(table, rows: list[tuple[str, str]]) -> None:
    """Rebuild a 2-column table's data rows from (col0, col1) tuples, cloning the
    first row as the formatting prototype. Header-less tables only."""
    tbl = table._tbl
    trs = tbl.findall(qn("a:tr"))
    if not trs:
        return
    proto = copy.deepcopy(trs[0])
    for tr in trs:
        tbl.remove(tr)
    for col0, col1 in rows:
        tr = copy.deepcopy(proto)
        tcs = tr.findall(qn("a:tc"))
        # python-pptx cell wrappers for set_cell_text
        from pptx.table import _Cell
        if len(tcs) >= 1:
            set_cell_text(_Cell(tcs[0], None), col0)
        if len(tcs) >= 2:
            set_cell_text(_Cell(tcs[1], None), col1)
        tbl.append(tr)
